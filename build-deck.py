#!/usr/bin/env python3
"""Build Meridian Health demo deck as .pptx (opens in Keynote)."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

prs = Presentation()
prs.slide_width = Emu(12192000)   # 1920px at 96dpi → ~13.33in but let's use standard 16:9
prs.slide_height = Emu(6858000)   # 1080px

NAVY = RGBColor(0x0A, 0x2E, 0x48)
TEAL = RGBColor(0x5D, 0xBE, 0xB7)
MINT = RGBColor(0x5D, 0xBE, 0xB7)
LIGHT_MINT = RGBColor(0x99, 0xD5, 0xD3)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xEC, 0xF2, 0xFD)
BODY_GRAY = RGBColor(0x64, 0x6E, 0x82)
MUTED = RGBColor(0x96, 0xA0, 0xB4)
OFFWHITE = RGBColor(0xEC, 0xF2, 0xFD)

BLANK = 6  # blank layout index

def add_textbox(slide, left, top, width, height, text, font_size=16, font_name="Calibri",
                color=NAVY, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.name = font_name
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return txBox, tf, p

def add_meridian_wordmark(slide, left, top, dark_bg=False):
    """Add 'meridian health' wordmark with split coloring."""
    txBox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(3500000), Emu(500000))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    
    run1 = p.add_run()
    run1.text = "meridian "
    run1.font.size = Pt(24)
    run1.font.name = "DM Sans"
    run1.font.bold = True
    run1.font.color.rgb = LIGHT_MINT if dark_bg else NAVY
    
    run2 = p.add_run()
    run2.text = "health"
    run2.font.size = Pt(24)
    run2.font.name = "DM Sans"
    run2.font.bold = False
    run2.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF, ) if dark_bg else TEAL
    # Note: 35% opacity white not possible in pptx, using full white for dark bg

def add_slide_number(slide, num, dark_bg=False):
    add_textbox(slide, 11500000, 6500000, 500000, 300000,
                f"{num:02d}", font_size=12, color=MUTED if not dark_bg else RGBColor(0x3C, 0x3C, 0x50))

def set_bg_color(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

# ============================================
# SLIDE 1: The story of your future
# ============================================
slide1 = prs.slides.add_slide(prs.slide_layouts[BLANK])
set_bg_color(slide1, LIGHT_BG)
add_meridian_wordmark(slide1, 400000, 250000)

# Main heading - need split color
txBox = slide1.shapes.add_textbox(Emu(1000000), Emu(800000), Emu(10192000), Emu(1800000))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER

run1 = p.add_run()
run1.text = "The story of your future is written\non top of "
run1.font.size = Pt(52)
run1.font.name = "Georgia"
run1.font.bold = True
run1.font.color.rgb = NAVY

run2 = p.add_run()
run2.text = "two curves."
run2.font.size = Pt(52)
run2.font.name = "Georgia"
run2.font.bold = True
run2.font.color.rgb = TEAL

# Placeholder boxes for charts
from pptx.util import Inches as In
from pptx.enum.shapes import MSO_SHAPE

# Left chart placeholder
left_box = slide1.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Emu(500000), Emu(2800000), Emu(5400000), Emu(3500000))
left_box.fill.solid()
left_box.fill.fore_color.rgb = WHITE
left_box.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
left_box.line.width = Pt(1)

add_textbox(slide1, 700000, 2950000, 4000000, 400000,
            "Quality of life", font_size=20, bold=True, color=NAVY)
add_textbox(slide1, 700000, 3300000, 4800000, 400000,
            "What your body can do at every age — from full capability to needing help.",
            font_size=13, color=BODY_GRAY)

# Right chart placeholder
right_box = slide1.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Emu(6300000), Emu(2800000), Emu(5400000), Emu(3500000))
right_box.fill.solid()
right_box.fill.fore_color.rgb = WHITE
right_box.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
right_box.line.width = Pt(1)

add_textbox(slide1, 6500000, 2950000, 4000000, 400000,
            "Probability of survival", font_size=20, bold=True, color=NAVY)
add_textbox(slide1, 6500000, 3300000, 4800000, 400000,
            "Your likelihood of being alive at each age, based on your health profile.",
            font_size=13, color=BODY_GRAY)

add_slide_number(slide1, 1)

# ============================================
# SLIDE 2: Blood, Behavior, DNA
# ============================================
slide2 = prs.slides.add_slide(prs.slide_layouts[BLANK])
set_bg_color(slide2, LIGHT_BG)
add_meridian_wordmark(slide2, 400000, 250000)

# Heading with bold keywords
txBox2 = slide2.shapes.add_textbox(Emu(1500000), Emu(1200000), Emu(9192000), Emu(2000000))
tf2 = txBox2.text_frame
tf2.word_wrap = True
p2 = tf2.paragraphs[0]
p2.alignment = PP_ALIGN.CENTER

for text, bold, color in [
    ("Those curves are shaped by\nyour ", True, NAVY),
    ("blood", True, NAVY),
    (", ", True, NAVY),
    ("behavior", True, NAVY),
    (",\nand ", True, NAVY),
    ("DNA", True, NAVY),
    (".", True, NAVY),
]:
    run = p2.add_run()
    run.text = text
    run.font.size = Pt(48)
    run.font.name = "Georgia"
    run.font.bold = bold
    run.font.color.rgb = color

# Three columns
col_width = 3200000
col_gap = 600000
col_start = 1000000
col_top = 3800000

cards = [
    ("🩸", "Blood", "Cholesterol, inflammation, hormones — the biomarkers that explain what your body is tracking right now."),
    ("🏃", "Behavior", "Fitness, sleep, training consistency — the daily choices that move your curve up or down."),
    ("🧬", "DNA", "Genetic variants that set your baseline — and help you know what matters most, what to test, and how to act."),
]

for i, (emoji, title, desc) in enumerate(cards):
    x = col_start + i * (col_width + col_gap)
    
    add_textbox(slide2, x, col_top - 400000, 500000, 400000, emoji, font_size=36)
    add_textbox(slide2, x, col_top, col_width, 350000, title, font_size=24, bold=True, color=NAVY)
    add_textbox(slide2, x, col_top + 400000, col_width, 600000, desc, font_size=14, color=BODY_GRAY)

add_slide_number(slide2, 2)

# ============================================
# SLIDE 3: Bend the curves
# ============================================
slide3 = prs.slides.add_slide(prs.slide_layouts[BLANK])
set_bg_color(slide3, LIGHT_BG)
add_meridian_wordmark(slide3, 400000, 250000)

txBox3 = slide3.shapes.add_textbox(Emu(1000000), Emu(1000000), Emu(10192000), Emu(1800000))
tf3 = txBox3.text_frame
tf3.word_wrap = True
p3 = tf3.paragraphs[0]
p3.alignment = PP_ALIGN.CENTER

run3a = p3.add_run()
run3a.text = "If you know what to do,\nyou can "
run3a.font.size = Pt(52)
run3a.font.name = "Georgia"
run3a.font.bold = True
run3a.font.color.rgb = NAVY

run3b = p3.add_run()
run3b.text = "bend the curves."
run3b.font.size = Pt(52)
run3b.font.name = "Georgia"
run3b.font.bold = True
run3b.font.color.rgb = TEAL

# Chart placeholders (same layout as slide 1)
left_box3 = slide3.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Emu(500000), Emu(2800000), Emu(5400000), Emu(3500000))
left_box3.fill.solid()
left_box3.fill.fore_color.rgb = WHITE
left_box3.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
left_box3.line.width = Pt(1)

right_box3 = slide3.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Emu(6300000), Emu(2800000), Emu(5400000), Emu(3500000))
right_box3.fill.solid()
right_box3.fill.fore_color.rgb = WHITE
right_box3.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
right_box3.line.width = Pt(1)

add_textbox(slide3, 700000, 2950000, 4000000, 400000,
            "Quality of life", font_size=20, bold=True, color=NAVY)
add_textbox(slide3, 700000, 3300000, 4800000, 400000,
            "The same curve — but with the right interventions, it stays higher for longer.",
            font_size=13, color=BODY_GRAY)

add_textbox(slide3, 6500000, 2950000, 4000000, 400000,
            "Probability of survival", font_size=20, bold=True, color=NAVY)
add_textbox(slide3, 6500000, 3300000, 4800000, 400000,
            "The gap between the lines is what knowing — and acting — is worth.",
            font_size=13, color=BODY_GRAY)

add_slide_number(slide3, 3)

# ============================================
# SLIDE 4: 30 more years
# ============================================
slide4 = prs.slides.add_slide(prs.slide_layouts[BLANK])
set_bg_color(slide4, NAVY)

add_textbox(slide4, 600000, 1500000, 5000000, 2200000,
            "30", font_size=180, font_name="Georgia", bold=True, color=MINT)

add_textbox(slide4, 600000, 3400000, 5000000, 500000,
            "more years", font_size=36, bold=True, color=OFFWHITE)

add_textbox(slide4, 600000, 3900000, 5000000, 600000,
            "of being anything you want.\nThat's the gap between average and best possible.",
            font_size=16, color=MUTED)

# Chart placeholder on right
right_box4 = slide4.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Emu(6300000), Emu(1500000), Emu(5400000), Emu(3500000))
right_box4.fill.solid()
right_box4.fill.fore_color.rgb = RGBColor(0x12, 0x3A, 0x55)
right_box4.line.color.rgb = RGBColor(0x1A, 0x45, 0x60)
right_box4.line.width = Pt(1)

add_meridian_wordmark(slide4, 400000, 6300000, dark_bg=True)
add_slide_number(slide4, 4, dark_bg=True)

# ============================================
# SLIDE 5: What if?
# ============================================
slide5 = prs.slides.add_slide(prs.slide_layouts[BLANK])
set_bg_color(slide5, LIGHT_BG)
add_meridian_wordmark(slide5, 400000, 250000)

add_textbox(slide5, 1000000, 1000000, 10192000, 1200000,
            "What if?", font_size=80, font_name="Georgia", bold=True,
            color=NAVY, alignment=PP_ALIGN.CENTER)

# Left pill (Meridian)
left_pill = slide5.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Emu(1500000), Emu(3200000), Emu(3800000), Emu(1400000))
left_pill.fill.solid()
left_pill.fill.fore_color.rgb = NAVY

# Meridian text in pill
mer_box = slide5.shapes.add_textbox(Emu(1700000), Emu(3400000), Emu(3400000), Emu(500000))
tf_mer = mer_box.text_frame
p_mer = tf_mer.paragraphs[0]
p_mer.alignment = PP_ALIGN.CENTER
r1 = p_mer.add_run()
r1.text = "meridian "
r1.font.size = Pt(28)
r1.font.name = "DM Sans"
r1.font.bold = True
r1.font.color.rgb = LIGHT_MINT
r2 = p_mer.add_run()
r2.text = "health"
r2.font.size = Pt(28)
r2.font.name = "DM Sans"
r2.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

add_textbox(slide5, 1700000, 3900000, 3400000, 300000,
            "HEALTH INTELLIGENCE", font_size=11, color=MUTED, alignment=PP_ALIGN.CENTER)

# × connector
add_textbox(slide5, 5350000, 3500000, 500000, 500000,
            "×", font_size=36, color=MUTED, alignment=PP_ALIGN.CENTER)

# Right pill (Colonial Life)
right_pill = slide5.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Emu(5900000), Emu(3200000), Emu(3800000), Emu(1400000))
right_pill.fill.solid()
right_pill.fill.fore_color.rgb = RGBColor(0x09, 0x2C, 0x48)

add_textbox(slide5, 6100000, 3500000, 3400000, 500000,
            "Colonial Life.", font_size=28, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

add_textbox(slide5, 6100000, 3900000, 3400000, 300000,
            "CRITICAL ILLNESS", font_size=11, color=MUTED, alignment=PP_ALIGN.CENTER)

# = Executive CI
add_textbox(slide5, 9900000, 3500000, 2200000, 500000,
            "= Executive CI", font_size=32, font_name="Georgia", bold=True, color=NAVY)

# Subtitle
add_textbox(slide5, 1500000, 4900000, 9192000, 400000,
            "A new category: critical illness insurance supercharged with continuous health intelligence.",
            font_size=15, color=BODY_GRAY, alignment=PP_ALIGN.CENTER)

add_slide_number(slide5, 5)

# ============================================
# SLIDE 6: What it could mean
# ============================================
slide6 = prs.slides.add_slide(prs.slide_layouts[BLANK])
set_bg_color(slide6, NAVY)
add_meridian_wordmark(slide6, 400000, 250000, dark_bg=True)

# Heading
txBox6 = slide6.shapes.add_textbox(Emu(1000000), Emu(800000), Emu(10192000), Emu(1000000))
tf6 = txBox6.text_frame
p6 = tf6.paragraphs[0]
p6.alignment = PP_ALIGN.CENTER
r6a = p6.add_run()
r6a.text = "What it could "
r6a.font.size = Pt(48)
r6a.font.name = "Georgia"
r6a.font.bold = True
r6a.font.color.rgb = OFFWHITE
r6b = p6.add_run()
r6b.text = "mean."
r6b.font.size = Pt(48)
r6b.font.name = "Georgia"
r6b.font.bold = True
r6b.font.color.rgb = MINT

# Three cards
card_width = 3200000
card_height = 2800000
card_gap = 400000
card_start = 700000
card_top = 2200000

cards6 = [
    ("2–3×", "Broker commissions",
     "Meridian layer on top of base CI gives brokers a reason to sell — and earn 2–3× the commission of the base product alone."),
    ("Take share", "Displace the incumbent",
     "A differentiated product gives brokers a reason to re-open every account. New story, new sale."),
    ("↓ Claims", "Increased margins",
     "If Meridian catches disease earlier and members act on it, the carrier's loss ratio improves. Better outcomes, lower claims, aligned incentives."),
]

for i, (headline, subtitle, desc) in enumerate(cards6):
    x = card_start + i * (card_width + card_gap)
    
    # Card background
    card_bg = slide6.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Emu(x), Emu(card_top), Emu(card_width), Emu(card_height))
    card_bg.fill.solid()
    card_bg.fill.fore_color.rgb = RGBColor(0x10, 0x38, 0x52)
    card_bg.line.color.rgb = RGBColor(0x1A, 0x45, 0x60)
    card_bg.line.width = Pt(1)
    
    add_textbox(slide6, x + 250000, card_top + 250000, card_width - 500000, 700000,
                headline, font_size=42, font_name="Georgia", bold=True, color=MINT)
    
    add_textbox(slide6, x + 250000, card_top + 900000, card_width - 500000, 350000,
                subtitle, font_size=18, bold=True, color=OFFWHITE)
    
    add_textbox(slide6, x + 250000, card_top + 1300000, card_width - 500000, 800000,
                desc, font_size=14, color=MUTED)

add_slide_number(slide6, 6, dark_bg=True)

# Save
out_path = os.path.expanduser("~/.openclaw/workspace/dashboard/meridian-demo.pptx")
prs.save(out_path)
print(f"Saved to {out_path}")

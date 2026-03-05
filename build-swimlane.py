#!/usr/bin/env python3
"""Build swim lane slide as .pptx for Keynote editing."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Emu(12192000)   # 16:9
prs.slide_height = Emu(6858000)

NAVY = RGBColor(0x11, 0x23, 0x4B)
TEAL = RGBColor(0x0D, 0x80, 0x86)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF5, 0xF9, 0xFF)
BODY_GRAY = RGBColor(0x54, 0x56, 0x6B)
MUTED = RGBColor(0x96, 0xA0, 0xB4)
BORDER_GRAY = RGBColor(0xE0, 0xE8, 0xF0)
CARD_BG = RGBColor(0xFF, 0xFF, 0xFF)
YOU_BG = RGBColor(0xF0, 0xF2, 0xF8)      # light navy tint for YOU lane
MER_BG = RGBColor(0xEE, 0xF8, 0xF8)      # light teal tint for MERIDIAN lane
TEAL_LIGHT = RGBColor(0xE0, 0xF5, 0xF4)

BLANK = 6

def emu(inches):
    return Emu(int(inches * 914400))

def add_textbox(slide, left, top, width, height, text="", font_size=16,
                font_name="Calibri", color=NAVY, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    if text:
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.name = font_name
        p.font.color.rgb = color
        p.font.bold = bold
        p.alignment = alignment
    return txBox

def add_rounded_rect(slide, left, top, width, height, fill_color=CARD_BG, 
                     line_color=None, line_width=Pt(1)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Emu(left), Emu(top), Emu(width), Emu(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    # Adjust corner radius
    shape.adjustments[0] = 0.1  # ~10% corner radius
    return shape

def add_rich_textbox(slide, left, top, width, height):
    """Return (shape, text_frame) for multi-run text."""
    txBox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    return txBox, tf

# ============================================
# SLIDE: One-time setup. Continuous progress.
# ============================================
slide = prs.slides.add_slide(prs.slide_layouts[BLANK])

# Background
bg = slide.background
bg.fill.solid()
bg.fill.fore_color.rgb = LIGHT_BG

# --- Wordmark ---
txBox, tf = add_rich_textbox(slide, 350000, 200000, 3500000, 450000)
p = tf.paragraphs[0]
r1 = p.add_run()
r1.text = "meridian "
r1.font.size = Pt(20)
r1.font.bold = True
r1.font.color.rgb = NAVY
r2 = p.add_run()
r2.text = "health"
r2.font.size = Pt(20)
r2.font.bold = False
r2.font.color.rgb = TEAL

# --- Headline ---
txBox, tf = add_rich_textbox(slide, 500000, 550000, 11000000, 900000)
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r1 = p.add_run()
r1.text = "One-time setup. "
r1.font.size = Pt(40)
r1.font.bold = True
r1.font.color.rgb = NAVY
r1.font.name = "DM Sans"
r2 = p.add_run()
r2.text = "Continuous progress."
r2.font.size = Pt(40)
r2.font.bold = True
r2.font.color.rgb = TEAL
r2.font.name = "DM Sans"

# --- Subtitle ---
add_textbox(slide, 500000, 1350000, 11000000, 400000,
            "You onboard once — then enter a cycle of intelligent progression.",
            font_size=14, color=BODY_GRAY, alignment=PP_ALIGN.CENTER)

# ============================================
# SWIM LANE GRID
# ============================================
# Layout constants
LANE_TOP = 1900000       # where the grid starts
COL_LABEL_W = 900000     # width for YOU/MERIDIAN labels
LANE_H = 1900000         # height per lane
DIVIDER_Y = LANE_TOP + LANE_H  # horizontal divider between lanes
TOTAL_H = LANE_H * 2

# Column positions (after label column)
GRID_LEFT = 900000
COL_W_SETUP = 1800000
COL_W_MONTHS = 3200000
COL_W_M6 = 1700000
COL_W_M712 = 3200000

COL1_L = GRID_LEFT                          # Setup
COL2_L = COL1_L + COL_W_SETUP + 80000      # Months 1-6
COL3_L = COL2_L + COL_W_MONTHS + 80000     # Month 6
COL4_L = COL3_L + COL_W_M6 + 80000        # Months 7-12

# --- Phase headers ---
header_y = LANE_TOP - 280000
for (x, w, label) in [
    (COL1_L, COL_W_SETUP, "SETUP"),
    (COL2_L, COL_W_MONTHS, "MONTHS 1–6"),
    (COL3_L, COL_W_M6, "MONTH 6"),
    (COL4_L, COL_W_M712, "MONTHS 7–12"),
]:
    add_textbox(slide, x, header_y, w, 250000, label,
                font_size=9, color=TEAL, bold=True, alignment=PP_ALIGN.CENTER)

# --- Lane labels ---
add_textbox(slide, 100000, LANE_TOP + 600000, COL_LABEL_W, 400000, "YOU",
            font_size=18, color=NAVY, bold=True, alignment=PP_ALIGN.CENTER,
            font_name="DM Sans")
add_textbox(slide, 100000, DIVIDER_Y + 600000, COL_LABEL_W, 400000, "MERIDIAN",
            font_size=16, color=TEAL, bold=True, alignment=PP_ALIGN.CENTER,
            font_name="DM Sans")

# --- Lane background tints ---
# YOU lane background
you_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Emu(GRID_LEFT), Emu(LANE_TOP), Emu(11000000 - GRID_LEFT), Emu(LANE_H))
you_bg.fill.solid()
you_bg.fill.fore_color.rgb = RGBColor(0xF0, 0xF2, 0xF7)
you_bg.line.fill.background()

# MERIDIAN lane background
mer_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Emu(GRID_LEFT), Emu(DIVIDER_Y), Emu(11000000 - GRID_LEFT), Emu(LANE_H))
mer_bg.fill.solid()
mer_bg.fill.fore_color.rgb = RGBColor(0xEE, 0xF7, 0xF7)
mer_bg.line.fill.background()

# --- Horizontal divider ---
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Emu(GRID_LEFT), Emu(DIVIDER_Y - 10000), Emu(11000000 - GRID_LEFT), Emu(20000))
line.fill.solid()
line.fill.fore_color.rgb = BORDER_GRAY
line.line.fill.background()

# --- Vertical label divider ---
vline = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Emu(GRID_LEFT - 15000), Emu(LANE_TOP), Emu(20000), Emu(TOTAL_H))
vline.fill.solid()
vline.fill.fore_color.rgb = BORDER_GRAY
vline.line.fill.background()

# ============================================
# COLUMN 1: SETUP (spans both lanes)
# ============================================
CARD_PAD = 100000
setup_card = add_rounded_rect(slide,
    COL1_L + CARD_PAD, LANE_TOP + CARD_PAD,
    COL_W_SETUP - CARD_PAD*2, TOTAL_H - CARD_PAD*2,
    fill_color=WHITE, line_color=TEAL, line_width=Pt(2))

# Onboarding text
add_textbox(slide, COL1_L + CARD_PAD + 80000, LANE_TOP + 250000,
            COL_W_SETUP - 300000, 350000,
            "🧬  🩸  ⌚", font_size=22, alignment=PP_ALIGN.CENTER)

add_textbox(slide, COL1_L + CARD_PAD + 80000, LANE_TOP + 650000,
            COL_W_SETUP - 300000, 300000,
            "Onboarding", font_size=15, color=NAVY, bold=True, 
            alignment=PP_ALIGN.CENTER, font_name="DM Sans")

txBox, tf = add_rich_textbox(slide, COL1_L + CARD_PAD + 50000, LANE_TOP + 1000000,
                              COL_W_SETUP - 250000, 1200000)
items = ["Genetics kit", "Blood panel", "Wearable sync", "Initial assessment"]
for i, item in enumerate(items):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = item
    p.font.size = Pt(10)
    p.font.color.rgb = BODY_GRAY
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(3)

# ============================================
# COLUMN 2: MONTHS 1-6
# ============================================
# YOU lane card
you_m16 = add_rounded_rect(slide,
    COL2_L + CARD_PAD, LANE_TOP + CARD_PAD,
    COL_W_MONTHS - CARD_PAD*2, LANE_H - CARD_PAD*2,
    fill_color=WHITE, line_color=BORDER_GRAY)

add_textbox(slide, COL2_L + 200000, LANE_TOP + 200000,
            COL_W_MONTHS - 400000, 300000,
            "Work your priorities", font_size=14, color=NAVY, bold=True,
            font_name="DM Sans")

# Priority pills
pill_y = LANE_TOP + 550000
pill_labels = ["🏋️ Train", "😴 Sleep", "💊 Meds", "🥗 Nutrition"]
pill_x = COL2_L + 200000
for label in pill_labels:
    pill = add_rounded_rect(slide, pill_x, pill_y, 650000, 280000,
                            fill_color=RGBColor(0xF0, 0xF2, 0xF7), line_color=None)
    add_textbox(slide, pill_x + 40000, pill_y + 30000, 570000, 220000,
                label, font_size=9, color=NAVY, bold=True, alignment=PP_ALIGN.CENTER)
    pill_x += 720000

add_textbox(slide, COL2_L + 200000, LANE_TOP + 1000000,
            COL_W_MONTHS - 400000, 600000,
            "Wear your device daily. Follow your 3 priorities. Live your life.",
            font_size=10, color=BODY_GRAY)

# MERIDIAN lane card
mer_m16 = add_rounded_rect(slide,
    COL2_L + CARD_PAD, DIVIDER_Y + CARD_PAD,
    COL_W_MONTHS - CARD_PAD*2, LANE_H - CARD_PAD*2,
    fill_color=WHITE, line_color=BORDER_GRAY)

add_textbox(slide, COL2_L + 200000, DIVIDER_Y + 200000,
            COL_W_MONTHS - 400000, 300000,
            "Working behind the scenes", font_size=14, color=TEAL, bold=True,
            font_name="DM Sans")

# Meridian pills
mpill_y = DIVIDER_Y + 550000
mpill_labels = ["📊 Monitor", "🔬 Research", "🎯 Adjust"]
mpill_x = COL2_L + 200000
for label in mpill_labels:
    pill = add_rounded_rect(slide, mpill_x, mpill_y, 800000, 280000,
                            fill_color=RGBColor(0xEE, 0xF7, 0xF7), line_color=None)
    add_textbox(slide, mpill_x + 40000, mpill_y + 30000, 720000, 220000,
                label, font_size=9, color=TEAL, bold=True, alignment=PP_ALIGN.CENTER)
    mpill_x += 870000

add_textbox(slide, COL2_L + 200000, DIVIDER_Y + 1000000,
            COL_W_MONTHS - 400000, 600000,
            "Track wearable data. Monitor progress. Scan relevant medical studies. Adjust coaching.",
            font_size=10, color=BODY_GRAY)

# ============================================
# COLUMN 3: MONTH 6 (spans both lanes)
# ============================================
m6_card = add_rounded_rect(slide,
    COL3_L + CARD_PAD, LANE_TOP + CARD_PAD,
    COL_W_M6 - CARD_PAD*2, TOTAL_H - CARD_PAD*2,
    fill_color=WHITE, line_color=TEAL, line_width=Pt(2))

add_textbox(slide, COL3_L + CARD_PAD + 50000, LANE_TOP + 500000,
            COL_W_M6 - 300000, 350000,
            "🩸", font_size=28, alignment=PP_ALIGN.CENTER)

add_textbox(slide, COL3_L + CARD_PAD + 50000, LANE_TOP + 1000000,
            COL_W_M6 - 300000, 400000,
            "Blood panel #2", font_size=15, color=NAVY, bold=True,
            alignment=PP_ALIGN.CENTER, font_name="DM Sans")

add_textbox(slide, COL3_L + CARD_PAD + 50000, LANE_TOP + 1350000,
            COL_W_M6 - 300000, 400000,
            "Updated\ngame plan", font_size=11, color=BODY_GRAY,
            alignment=PP_ALIGN.CENTER)

# ============================================
# COLUMN 4: MONTHS 7-12
# ============================================
# YOU lane card
you_m712 = add_rounded_rect(slide,
    COL4_L + CARD_PAD, LANE_TOP + CARD_PAD,
    COL_W_M712 - CARD_PAD*2, LANE_H - CARD_PAD*2,
    fill_color=WHITE, line_color=BORDER_GRAY)

add_textbox(slide, COL4_L + 200000, LANE_TOP + 200000,
            COL_W_M712 - 400000, 300000,
            "Refined priorities", font_size=14, color=NAVY, bold=True,
            font_name="DM Sans")

add_textbox(slide, COL4_L + 200000, LANE_TOP + 550000,
            COL_W_M712 - 400000, 600000,
            "Updated based on your results. What's working stays. What's not gets adjusted.",
            font_size=10, color=BODY_GRAY)

# "repeats →" indicator
add_textbox(slide, COL4_L + 200000, LANE_TOP + 1300000,
            COL_W_M712 - 400000, 250000,
            "repeats →", font_size=9, color=MUTED, bold=True, alignment=PP_ALIGN.RIGHT)

# MERIDIAN lane card
mer_m712 = add_rounded_rect(slide,
    COL4_L + CARD_PAD, DIVIDER_Y + CARD_PAD,
    COL_W_M712 - CARD_PAD*2, LANE_H - CARD_PAD*2,
    fill_color=WHITE, line_color=BORDER_GRAY)

add_textbox(slide, COL4_L + 200000, DIVIDER_Y + 200000,
            COL_W_M712 - 400000, 300000,
            "Smarter every cycle", font_size=14, color=TEAL, bold=True,
            font_name="DM Sans")

add_textbox(slide, COL4_L + 200000, DIVIDER_Y + 550000,
            COL_W_M712 - 400000, 600000,
            "More data = better insights. AI recalibrates risk. New research integrated. Your picture sharpens.",
            font_size=10, color=BODY_GRAY)

add_textbox(slide, COL4_L + 200000, DIVIDER_Y + 1300000,
            COL_W_M712 - 400000, 250000,
            "repeats →", font_size=9, color=MUTED, bold=True, alignment=PP_ALIGN.RIGHT)

# --- Bottom summary pill ---
pill_w = 7000000
pill_l = (12192000 - pill_w) // 2
summary_bg = add_rounded_rect(slide, pill_l, LANE_TOP + TOTAL_H + 200000,
                               pill_w, 350000,
                               fill_color=TEAL_LIGHT, line_color=None)
add_textbox(slide, pill_l + 100000, LANE_TOP + TOTAL_H + 230000,
            pill_w - 200000, 300000,
            "📋 Initial game plan presented → you start working → Meridian keeps optimizing",
            font_size=11, color=TEAL, bold=True, alignment=PP_ALIGN.CENTER)

# --- Slide number ---
add_textbox(slide, 11500000, 6500000, 500000, 300000,
            "05", font_size=11, color=MUTED)

# Save
out = os.path.join(os.path.dirname(__file__), "meridian-swimlane.pptx")
prs.save(out)
print(f"✅ Saved: {out}")
